from pathlib import Path
import csv
import io
import re
from xml.sax.saxutils import escape

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table as ExcelTable
from openpyxl.worksheet.table import TableStyleInfo

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
)


# =========================================================
# OUTPUT CONFIGURATION
# =========================================================

OUTPUT_DIR = Path("generated_documents")


# =========================================================
# DIRECTORY
# =========================================================

def ensure_output_directory():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# =========================================================
# COMMON HELPERS
# =========================================================

def clean_inline_markdown(text: str) -> str:
    text = str(text)

    text = re.sub(
        r"!\[([^\]]*)\]\([^)]+\)",
        r"\1",
        text,
    )

    text = re.sub(
        r"\[([^\]]+)\]\([^)]+\)",
        r"\1",
        text,
    )

    text = re.sub(
        r"\*\*(.*?)\*\*",
        r"\1",
        text,
    )

    text = re.sub(
        r"(?<!\*)\*([^*]+?)\*(?!\*)",
        r"\1",
        text,
    )

    text = re.sub(
        r"`([^`]+)`",
        r"\1",
        text,
    )

    return text.strip()


def parse_content(content: str) -> list[dict]:
    """
    Parse Markdown-like LLM output into simple blocks.
    """

    lines = str(content).replace("\r\n", "\n").split("\n")

    blocks = []
    paragraph_buffer = []

    def flush_paragraph():
        if paragraph_buffer:
            text = " ".join(
                item.strip()
                for item in paragraph_buffer
                if item.strip()
            )

            if text:
                blocks.append({
                    "type": "paragraph",
                    "text": text,
                })

            paragraph_buffer.clear()

    index = 0

    while index < len(lines):
        line = lines[index].strip()

        if not line:
            flush_paragraph()
            index += 1
            continue

        if line.startswith("# "):
            flush_paragraph()
            blocks.append({
                "type": "h1",
                "text": line[2:].strip(),
            })
            index += 1
            continue

        if line.startswith("## "):
            flush_paragraph()
            blocks.append({
                "type": "h2",
                "text": line[3:].strip(),
            })
            index += 1
            continue

        if line.startswith("### "):
            flush_paragraph()
            blocks.append({
                "type": "h3",
                "text": line[4:].strip(),
            })
            index += 1
            continue

        if re.match(r"^[-*_]{3,}$", line):
            flush_paragraph()
            blocks.append({
                "type": "rule",
                "text": "",
            })
            index += 1
            continue

        bullet_match = re.match(r"^[-*+]\s+(.*)", line)

        if bullet_match:
            flush_paragraph()
            blocks.append({
                "type": "bullet",
                "text": bullet_match.group(1),
            })
            index += 1
            continue

        numbered_match = re.match(r"^(\d+)\.\s+(.*)", line)

        if numbered_match:
            flush_paragraph()
            blocks.append({
                "type": "number",
                "number": numbered_match.group(1),
                "text": numbered_match.group(2),
            })
            index += 1
            continue

        if line.startswith(">"):
            flush_paragraph()
            blocks.append({
                "type": "quote",
                "text": line[1:].strip(),
            })
            index += 1
            continue

        paragraph_buffer.append(line)
        index += 1

    flush_paragraph()

    return blocks


def markdown_to_reportlab(text: str) -> str:
    text = clean_inline_markdown(text)
    return escape(text)


# =========================================================
# TXT
# =========================================================

def create_txt(
    filename: str,
    content: str,
) -> str:
    ensure_output_directory()

    file_path = OUTPUT_DIR / filename

    file_path.write_text(
        str(content),
        encoding="utf-8",
    )

    return str(file_path)


# =========================================================
# DOCX HELPERS
# =========================================================

def set_cell_shading(cell, fill="F2F2F2"):
    tc_pr = cell._tc.get_or_add_tcPr()

    shd = OxmlElement("w:shd")

    shd.set(
        qn("w:fill"),
        fill,
    )

    tc_pr.append(shd)


def set_cell_borders(cell):
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()

    tc_borders = tc_pr.first_child_found_in("w:tcBorders")

    if tc_borders is None:
        tc_borders = OxmlElement("w:tcBorders")
        tc_pr.append(tc_borders)

    for edge in (
        "top",
        "left",
        "bottom",
        "right",
        "insideH",
        "insideV",
    ):
        tag = "w:" + edge

        element = tc_borders.find(qn(tag))

        if element is None:
            element = OxmlElement(tag)
            tc_borders.append(element)

        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), "4")
        element.set(qn("w:color"), "D9D9D9")


# =========================================================
# DOCX
# =========================================================

def create_docx(
    filename: str,
    content: str,
) -> str:
    ensure_output_directory()

    file_path = OUTPUT_DIR / filename

    document = Document()

    section = document.sections[0]

    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    styles = document.styles

    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)

    blocks = parse_content(content)

    title_used = False

    if blocks and blocks[0]["type"] == "h1":
        title = document.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        run = title.add_run(
            clean_inline_markdown(
                blocks[0]["text"]
            )
        )

        run.bold = True
        run.font.size = Pt(22)
        run.font.name = "Aptos Display"

        subtitle = document.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        subtitle_run = subtitle.add_run(
            "Generated by CHORUS"
        )

        subtitle_run.font.size = Pt(9)
        subtitle_run.italic = True

        title_used = True

    start_index = 1 if title_used else 0

    for block in blocks[start_index:]:
        block_type = block["type"]

        text = clean_inline_markdown(
            block.get("text", "")
        )

        if block_type == "h1":
            paragraph = document.add_paragraph()

            paragraph.paragraph_format.space_before = Pt(16)
            paragraph.paragraph_format.space_after = Pt(8)

            run = paragraph.add_run(text)
            run.bold = True
            run.font.size = Pt(16)

        elif block_type == "h2":
            paragraph = document.add_paragraph()

            paragraph.paragraph_format.space_before = Pt(13)
            paragraph.paragraph_format.space_after = Pt(6)

            run = paragraph.add_run(text)
            run.bold = True
            run.font.size = Pt(13)

        elif block_type == "h3":
            paragraph = document.add_paragraph()

            paragraph.paragraph_format.space_before = Pt(10)
            paragraph.paragraph_format.space_after = Pt(5)

            run = paragraph.add_run(text)
            run.bold = True
            run.font.size = Pt(11)

        elif block_type == "bullet":
            paragraph = document.add_paragraph(
                style="List Bullet"
            )

            paragraph.paragraph_format.space_after = Pt(4)
            paragraph.add_run(text)

        elif block_type == "number":
            paragraph = document.add_paragraph(
                style="List Number"
            )

            paragraph.paragraph_format.space_after = Pt(4)
            paragraph.add_run(text)

        elif block_type == "quote":
            paragraph = document.add_paragraph()

            paragraph.paragraph_format.left_indent = Inches(0.3)
            paragraph.paragraph_format.right_indent = Inches(0.2)
            paragraph.paragraph_format.space_before = Pt(5)
            paragraph.paragraph_format.space_after = Pt(8)

            run = paragraph.add_run(text)
            run.italic = True

        elif block_type == "rule":
            paragraph = document.add_paragraph()

            paragraph.paragraph_format.space_before = Pt(4)
            paragraph.paragraph_format.space_after = Pt(8)

            run = paragraph.add_run(
                "────────────────────────────"
            )

            run.font.size = Pt(8)

        else:
            paragraph = document.add_paragraph()

            paragraph.paragraph_format.space_after = Pt(8)
            paragraph.paragraph_format.line_spacing = 1.15

            paragraph.add_run(text)

    footer = section.footer.paragraphs[0]

    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

    footer_run = footer.add_run(
        "Generated by CHORUS"
    )

    footer_run.font.size = Pt(8)
    footer_run.italic = True

    document.save(file_path)

    return str(file_path)


# =========================================================
# PDF
# =========================================================

def create_pdf(
    filename: str,
    content: str,
) -> str:
    """
    Create a professionally formatted PDF.

    IMPORTANT:
    This function is defined before create_document so the
    dispatcher can always resolve create_pdf at runtime.
    """

    ensure_output_directory()

    file_path = OUTPUT_DIR / filename

    document = SimpleDocTemplate(
        str(file_path),
        pagesize=A4,
        rightMargin=20 * mm,
        leftMargin=20 * mm,
        topMargin=22 * mm,
        bottomMargin=20 * mm,
        title=filename,
        author="CHORUS",
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "ChorusTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=22,
        leading=27,
        alignment=TA_CENTER,
        spaceAfter=10,
    )

    subtitle_style = ParagraphStyle(
        "ChorusSubtitle",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=9,
        leading=12,
        alignment=TA_CENTER,
        textColor=colors.grey,
        spaceAfter=20,
    )

    h1_style = ParagraphStyle(
        "ChorusH1",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=16,
        leading=20,
        spaceBefore=14,
        spaceAfter=8,
        keepWithNext=True,
    )

    h2_style = ParagraphStyle(
        "ChorusH2",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=17,
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True,
    )

    h3_style = ParagraphStyle(
        "ChorusH3",
        parent=styles["Heading3"],
        fontName="Helvetica-Bold",
        fontSize=11,
        leading=14,
        spaceBefore=9,
        spaceAfter=5,
        keepWithNext=True,
    )

    body_style = ParagraphStyle(
        "ChorusBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        spaceAfter=8,
    )

    bullet_style = ParagraphStyle(
        "ChorusBullet",
        parent=body_style,
        leftIndent=14,
        firstLineIndent=-8,
        spaceAfter=5,
    )

    quote_style = ParagraphStyle(
        "ChorusQuote",
        parent=body_style,
        leftIndent=16,
        rightIndent=12,
        fontName="Helvetica-Oblique",
        textColor=colors.grey,
    )

    elements = []

    blocks = parse_content(content)

    title_used = False

    if blocks and blocks[0]["type"] == "h1":
        elements.append(
            Paragraph(
                markdown_to_reportlab(
                    blocks[0]["text"]
                ),
                title_style,
            )
        )

        elements.append(
            Paragraph(
                "Generated by CHORUS",
                subtitle_style,
            )
        )

        title_used = True

    start_index = 1 if title_used else 0

    for block in blocks[start_index:]:
        block_type = block["type"]

        text = block.get("text", "")

        if block_type == "h1":
            elements.append(
                Paragraph(
                    markdown_to_reportlab(text),
                    h1_style,
                )
            )

        elif block_type == "h2":
            elements.append(
                Paragraph(
                    markdown_to_reportlab(text),
                    h2_style,
                )
            )

        elif block_type == "h3":
            elements.append(
                Paragraph(
                    markdown_to_reportlab(text),
                    h3_style,
                )
            )

        elif block_type == "bullet":
            elements.append(
                Paragraph(
                    "• "
                    + markdown_to_reportlab(text),
                    bullet_style,
                )
            )

        elif block_type == "number":
            elements.append(
                Paragraph(
                    f"{block['number']}. "
                    + markdown_to_reportlab(text),
                    bullet_style,
                )
            )

        elif block_type == "quote":
            elements.append(
                Paragraph(
                    markdown_to_reportlab(text),
                    quote_style,
                )
            )

        elif block_type == "rule":
            rule = Table(
                [[""]],
                colWidths=[170 * mm],
                rowHeights=[0.5 * mm],
            )

            rule.setStyle(
                TableStyle([
                    (
                        "BACKGROUND",
                        (0, 0),
                        (-1, -1),
                        colors.lightgrey,
                    )
                ])
            )

            elements.append(Spacer(1, 5))
            elements.append(rule)
            elements.append(Spacer(1, 8))

        else:
            elements.append(
                Paragraph(
                    markdown_to_reportlab(text),
                    body_style,
                )
            )

    if not elements:
        elements.append(
            Paragraph(
                "No content was generated.",
                body_style,
            )
        )

    def add_page_decoration(canvas, doc):
        canvas.saveState()

        width, height = A4

        canvas.setFont(
            "Helvetica",
            8,
        )

        canvas.setFillColor(
            colors.grey
        )

        canvas.drawString(
            20 * mm,
            height - 12 * mm,
            "CHORUS",
        )

        canvas.drawRightString(
            width - 20 * mm,
            10 * mm,
            f"Page {doc.page}",
        )

        canvas.restoreState()

    document.build(
        elements,
        onFirstPage=add_page_decoration,
        onLaterPages=add_page_decoration,
    )

    return str(file_path)


# =========================================================
# PPTX
# =========================================================

# CHORUS presentation theme.
PPT_BG = "F7F8FC"
PPT_TEXT = "20242E"
PPT_MUTED = "667085"
PPT_ACCENT = "5B5BD6"
PPT_ACCENT_DARK = "3F3FA8"
PPT_ACCENT_LIGHT = "E9E9FF"
PPT_WHITE = "FFFFFF"
PPT_BORDER = "E1E4EC"
PPT_DARK = "171927"


def _ppt_rgb(value):
    from pptx.dml.color import RGBColor

    value = value.replace("#", "")

    return RGBColor(
        int(value[0:2], 16),
        int(value[2:4], 16),
        int(value[4:6], 16),
    )


def _ppt_background(slide, color=PPT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _ppt_rgb(color)


def _ppt_add_box(
    slide,
    left,
    top,
    width,
    height,
    fill=PPT_WHITE,
    line=PPT_BORDER,
    rounded=True,
):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE
        if rounded
        else MSO_SHAPE.RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = _ppt_rgb(fill)

    if line:
        shape.line.color.rgb = _ppt_rgb(line)
        shape.line.width = PptPt(0.8)
    else:
        shape.line.fill.background()

    return shape


def _ppt_add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=16,
    bold=False,
    color=PPT_TEXT,
    font_name="Aptos",
    align=None,
):
    box = slide.shapes.add_textbox(
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    frame.margin_left = PptPt(1)
    frame.margin_right = PptPt(1)
    frame.margin_top = PptPt(1)
    frame.margin_bottom = PptPt(1)

    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.font.name = font_name
    paragraph.font.size = PptPt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _ppt_rgb(color)

    if align is not None:
        paragraph.alignment = align

    return box


def _ppt_add_accent_bar(slide):
    _ppt_add_box(
        slide,
        0,
        0,
        13.333,
        0.08,
        fill=PPT_ACCENT,
        line=None,
        rounded=False,
    )


def _ppt_add_footer(slide, slide_number):
    _ppt_add_text(
        slide,
        "CHORUS",
        0.65,
        7.03,
        1.3,
        0.22,
        font_size=8,
        bold=True,
        color=PPT_MUTED,
    )

    _ppt_add_text(
        slide,
        f"{slide_number:02d}",
        12.15,
        7.03,
        0.5,
        0.22,
        font_size=8,
        bold=True,
        color=PPT_MUTED,
        align=2,
    )


def _ppt_add_header(slide, title, subtitle=None):
    _ppt_add_text(
        slide,
        "CHORUS",
        0.7,
        0.42,
        2.2,
        0.22,
        font_size=8,
        bold=True,
        color=PPT_ACCENT,
    )

    _ppt_add_text(
        slide,
        title,
        0.7,
        0.78,
        11.8,
        0.7,
        font_size=27,
        bold=True,
        color=PPT_TEXT,
        font_name="Aptos Display",
    )

    if subtitle:
        _ppt_add_text(
            slide,
            subtitle,
            0.72,
            1.43,
            11.5,
            0.42,
            font_size=11,
            color=PPT_MUTED,
        )


def _ppt_add_bullet(
    slide,
    text,
    left,
    top,
    width,
    number=None,
):
    marker = str(number) if number is not None else "•"

    _ppt_add_box(
        slide,
        left,
        top + 0.06,
        0.28,
        0.28,
        fill=PPT_ACCENT_LIGHT,
        line=None,
        rounded=True,
    )

    _ppt_add_text(
        slide,
        marker,
        left,
        top + 0.055,
        0.28,
        0.28,
        font_size=8,
        bold=True,
        color=PPT_ACCENT,
        align=2,
    )

    _ppt_add_text(
        slide,
        text,
        left + 0.43,
        top,
        width - 0.43,
        0.65,
        font_size=14,
        color=PPT_TEXT,
    )


def _ppt_add_card(
    slide,
    label,
    body,
    left,
    top,
    width,
    height,
):
    _ppt_add_box(
        slide,
        left,
        top,
        width,
        height,
        fill=PPT_WHITE,
        line=PPT_BORDER,
        rounded=True,
    )

    _ppt_add_text(
        slide,
        label,
        left + 0.25,
        top + 0.22,
        width - 0.5,
        0.32,
        font_size=10,
        bold=True,
        color=PPT_ACCENT,
    )

    _ppt_add_text(
        slide,
        body,
        left + 0.25,
        top + 0.72,
        width - 0.5,
        height - 0.95,
        font_size=13,
        color=PPT_TEXT,
    )


def _ppt_add_title_slide(presentation, title, subtitle=None):
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    _ppt_background(slide, PPT_DARK)

    # Decorative CHORUS geometry.
    _ppt_add_box(
        slide,
        10.85,
        -0.35,
        2.9,
        2.9,
        fill=PPT_ACCENT_DARK,
        line=None,
        rounded=True,
    )

    _ppt_add_box(
        slide,
        11.75,
        5.1,
        1.85,
        1.85,
        fill=PPT_ACCENT,
        line=None,
        rounded=True,
    )

    _ppt_add_text(
        slide,
        "CHORUS",
        0.85,
        0.72,
        2.2,
        0.35,
        font_size=11,
        bold=True,
        color="BFC0FF",
    )

    _ppt_add_text(
        slide,
        title,
        0.85,
        2.0,
        10.2,
        1.5,
        font_size=35,
        bold=True,
        color=PPT_WHITE,
        font_name="Aptos Display",
    )

    if subtitle:
        _ppt_add_text(
            slide,
            subtitle,
            0.9,
            3.7,
            8.8,
            0.8,
            font_size=15,
            color="C7CBD8",
        )

    _ppt_add_text(
        slide,
        "Generated by CHORUS",
        0.9,
        6.45,
        3.0,
        0.3,
        font_size=9,
        color="9297A8",
    )

    return slide


def _ppt_build_sections(blocks):
    title = None
    start_index = 0

    if blocks and blocks[0]["type"] == "h1":
        title = clean_inline_markdown(
            blocks[0]["text"]
        )
        start_index = 1

    sections = []
    current = None

    for block in blocks[start_index:]:
        if block["type"] == "h2":
            if current:
                sections.append(current)

            current = {
                "title": clean_inline_markdown(
                    block["text"]
                ),
                "blocks": [],
            }

        else:
            if current is None:
                current = {
                    "title": title or "CHORUS",
                    "blocks": [],
                }

            current["blocks"].append(block)

    if current:
        sections.append(current)

    return title, sections


def _ppt_add_content_slide(
    presentation,
    section,
    slide_number,
):
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    _ppt_background(slide)
    _ppt_add_accent_bar(slide)
    _ppt_add_header(
        slide,
        section["title"],
    )

    blocks = section["blocks"]

    bullets = [
        block
        for block in blocks
        if block["type"] in ("bullet", "number")
    ]

    paragraphs = [
        block
        for block in blocks
        if block["type"] == "paragraph"
    ]

    headings = [
        block
        for block in blocks
        if block["type"] == "h3"
    ]

    quotes = [
        block
        for block in blocks
        if block["type"] == "quote"
    ]

    # -----------------------------------------------------
    # VISUAL CARD GRID
    # -----------------------------------------------------

    if (
        3 <= len(bullets) <= 6
        and not paragraphs
        and not headings
        and not quotes
    ):
        columns = 2 if len(bullets) <= 4 else 3
        gap = 0.3
        left = 0.7
        usable_width = 11.95
        card_width = (
            usable_width
            - gap * (columns - 1)
        ) / columns

        if columns == 2:
            card_height = 2.05
        else:
            card_height = 4.45

        for index, block in enumerate(bullets):
            row = index // columns
            column = index % columns

            _ppt_add_card(
                slide,
                f"{index + 1:02d}",
                clean_inline_markdown(
                    block["text"]
                ),
                left + column * (
                    card_width + gap
                ),
                1.95 + row * (
                    card_height + gap
                ),
                card_width,
                card_height,
            )

        _ppt_add_footer(
            slide,
            slide_number,
        )

        return slide

    # -----------------------------------------------------
    # TWO-COLUMN CONTENT
    # -----------------------------------------------------

    if bullets and paragraphs:
        _ppt_add_text(
            slide,
            "KEY POINTS",
            0.75,
            1.95,
            4.5,
            0.3,
            font_size=9,
            bold=True,
            color=PPT_ACCENT,
        )

        y = 2.35

        for block in bullets[:5]:
            _ppt_add_bullet(
                slide,
                clean_inline_markdown(
                    block["text"]
                ),
                0.75,
                y,
                5.15,
                (
                    block.get("number")
                    if block["type"] == "number"
                    else None
                ),
            )
            y += 0.82

        _ppt_add_card(
            slide,
            "CONTEXT",
            "\n\n".join(
                clean_inline_markdown(
                    paragraph["text"]
                )
                for paragraph in paragraphs[:2]
            ),
            6.35,
            1.95,
            5.75,
            4.35,
        )

        _ppt_add_footer(
            slide,
            slide_number,
        )

        return slide

    # -----------------------------------------------------
    # INSIGHT / QUOTE
    # -----------------------------------------------------

    if quotes:
        _ppt_add_card(
            slide,
            "KEY TAKEAWAY",
            clean_inline_markdown(
                quotes[0]["text"]
            ),
            0.8,
            2.0,
            11.7,
            2.1,
        )

        y = 4.45

        for block in (
            paragraphs + bullets
        )[:3]:
            _ppt_add_bullet(
                slide,
                clean_inline_markdown(
                    block["text"]
                ),
                0.95,
                y,
                10.9,
            )
            y += 0.7

        _ppt_add_footer(
            slide,
            slide_number,
        )

        return slide

    # -----------------------------------------------------
    # STANDARD CONTENT
    # -----------------------------------------------------

    y = 1.95

    content_blocks = (
        headings
        + paragraphs
        + bullets
    )

    for block in content_blocks[:8]:
        text = clean_inline_markdown(
            block.get("text", "")
        )

        if block["type"] == "h3":
            _ppt_add_text(
                slide,
                text,
                0.8,
                y,
                11.5,
                0.42,
                font_size=16,
                bold=True,
            )
            y += 0.58

        elif block["type"] in (
            "bullet",
            "number",
        ):
            _ppt_add_bullet(
                slide,
                text,
                0.8,
                y,
                11.5,
                (
                    block.get("number")
                    if block["type"] == "number"
                    else None
                ),
            )
            y += 0.76

        else:
            _ppt_add_text(
                slide,
                text,
                0.8,
                y,
                11.5,
                0.85,
                font_size=15,
            )
            y += 0.95

    if not content_blocks:
        _ppt_add_card(
            slide,
            "CHORUS",
            "No additional content was generated.",
            0.8,
            2.0,
            11.7,
            2.1,
        )

    _ppt_add_footer(
        slide,
        slide_number,
    )

    return slide


def create_pptx(
    filename: str,
    content: str,
) -> str:
    """
    Create a polished CHORUS PowerPoint.

    The LLM provides the document content.
    This generator provides the presentation design,
    hierarchy, spacing, branding and visual layouts.
    """

    ensure_output_directory()

    file_path = OUTPUT_DIR / filename

    presentation = Presentation()

    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)

    blocks = parse_content(content)

    title, sections = _ppt_build_sections(
        blocks
    )

    title = title or (
        sections[0]["title"]
        if sections
        else "CHORUS Presentation"
    )

    # First paragraph of the first section is used as a
    # lightweight subtitle on the title slide.
    subtitle = None

    if sections:
        for block in sections[0]["blocks"]:
            if block["type"] == "paragraph":
                subtitle = clean_inline_markdown(
                    block["text"]
                )
                break

    _ppt_add_title_slide(
        presentation,
        title,
        subtitle,
    )

    slide_number = 2

    for section in sections:
        if not section["blocks"]:
            continue

        _ppt_add_content_slide(
            presentation,
            section,
            slide_number,
        )

        slide_number += 1

    # Title-only requests still receive a useful second slide.
    if len(presentation.slides) == 1:
        _ppt_add_content_slide(
            presentation,
            {
                "title": "Overview",
                "blocks": [{
                    "type": "paragraph",
                    "text": (
                        "Presentation generated by CHORUS."
                    ),
                }],
            },
            2,
        )

    presentation.save(
        file_path
    )

    return str(file_path)


# =========================================================
# XLSX HELPERS
# =========================================================

def _content_to_rows(content) -> list[list]:
    """
    Convert common LLM table/CSV output into Excel rows.
    """

    if isinstance(content, list):
        if all(isinstance(row, (list, tuple)) for row in content):
            return [list(row) for row in content]

        return [[str(item)] for item in content]

    text = str(content).strip()

    if not text:
        return [["CHORUS"]]

    lines = [
        line.strip()
        for line in text.replace("\r\n", "\n").split("\n")
        if line.strip()
    ]

    # Markdown table
    markdown_rows = []

    for line in lines:
        if "|" not in line:
            continue

        stripped = line.strip().strip("|")

        cells = [
            cell.strip()
            for cell in stripped.split("|")
        ]

        if not cells:
            continue

        if all(
            re.fullmatch(r":?-{3,}:?", cell or "")
            for cell in cells
        ):
            continue

        markdown_rows.append(cells)

    if len(markdown_rows) >= 2:
        return markdown_rows

    # CSV / comma-separated
    try:
        reader = csv.reader(io.StringIO(text))
        rows = [
            [cell.strip() for cell in row]
            for row in reader
            if row
        ]

        if len(rows) >= 2 and any(
            len(row) > 1
            for row in rows
        ):
            return rows

    except Exception:
        pass

    # Fallback: preserve the generated content line-by-line.
    return [["CHORUS DOCUMENT CONTENT"]] + [
        [line]
        for line in lines
    ]


# =========================================================
# XLSX
# =========================================================

def create_xlsx(
    filename: str,
    rows,
) -> str:
    """
    Create a professionally formatted Excel spreadsheet.

    Accepts either:
    - list[list]
    - Markdown table text
    - CSV text
    - ordinary generated text
    """

    ensure_output_directory()

    file_path = OUTPUT_DIR / filename

    workbook = Workbook()

    worksheet = workbook.active
    worksheet.title = "CHORUS Data"

    rows = _content_to_rows(rows)

    for row in rows:
        worksheet.append(row)

    # Header
    if worksheet.max_row >= 1:
        for cell in worksheet[1]:
            cell.font = Font(
                bold=True
            )
            cell.alignment = Alignment(
                horizontal="center",
                vertical="center",
            )

    # General formatting
    thin = Side(
        style="thin",
        color="D9D9D9",
    )

    for row in worksheet.iter_rows():
        for cell in row:
            cell.border = Border(
                top=thin,
                bottom=thin,
                left=thin,
                right=thin,
            )
            cell.alignment = Alignment(
                vertical="top",
                wrap_text=True,
            )

    # Column widths
    for column_cells in worksheet.columns:
        column_letter = get_column_letter(
            column_cells[0].column
        )

        max_length = 0

        for cell in column_cells:
            value = (
                ""
                if cell.value is None
                else str(cell.value)
            )

            max_length = max(
                max_length,
                len(value),
            )

        worksheet.column_dimensions[
            column_letter
        ].width = min(
            max(max_length + 2, 12),
            40,
        )

    worksheet.freeze_panes = "A2"

    # Add an Excel table when there is actual tabular data.
    if (
        worksheet.max_row >= 2
        and worksheet.max_column >= 1
    ):
        last_column = get_column_letter(
            worksheet.max_column
        )

        table_ref = (
            f"A1:{last_column}{worksheet.max_row}"
        )

        table = ExcelTable(
            displayName="ChorusData",
            ref=table_ref,
        )

        style = TableStyleInfo(
            name="TableStyleMedium2",
            showFirstColumn=False,
            showLastColumn=False,
            showRowStripes=True,
            showColumnStripes=False,
        )

        table.tableStyleInfo = style

        worksheet.add_table(table)

    workbook.save(file_path)

    return str(file_path)


# =========================================================
# DOCUMENT CREATION
# =========================================================

def create_document(
    filename: str,
    content,
    file_type: str,
) -> str:
    """
    Central dispatcher for CHORUS document generation.

    Supported:
    - txt
    - docx
    - pdf
    - pptx
    - xlsx

    Each generator is defined in this module, including
    create_pdf. This avoids the runtime NameError that occurred
    when the PDF dispatcher could not resolve create_pdf.
    """

    file_type = (
        str(file_type)
        .lower()
        .replace(".", "")
        .strip()
    )

    if file_type == "txt":
        return create_txt(
            filename,
            str(content),
        )

    if file_type == "docx":
        return create_docx(
            filename,
            str(content),
        )

    if file_type == "pdf":
        return create_pdf(
            filename,
            str(content),
        )

    if file_type == "pptx":
        return create_pptx(
            filename,
            str(content),
        )

    if file_type == "xlsx":
        return create_xlsx(
            filename,
            content,
        )

    raise ValueError(
        f"Unsupported document type: {file_type}"
    )
