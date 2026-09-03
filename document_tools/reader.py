from pathlib import Path

import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def read_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""

    reader = PdfReader(file_path)

    text = []

    for page in reader.pages:
        page_text = page.extract_text()

        if page_text:
            text.append(page_text)

    return "\n".join(text).strip()


def read_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""

    document = Document(file_path)

    paragraphs = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    return "\n".join(paragraphs).strip()


def read_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint presentation."""

    presentation = Presentation(file_path)

    text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)

    return "\n".join(text).strip()


def read_csv(file_path: str) -> str:
    """Read a CSV file and return its contents."""

    dataframe = pd.read_csv(file_path)

    return dataframe.to_string(index=False)


def read_xlsx(file_path: str) -> str:
    """Read an Excel file and return its contents."""

    dataframe = pd.read_excel(file_path)

    return dataframe.to_string(index=False)


def read_txt(file_path: str) -> str:
    """Read a text file."""

    return Path(file_path).read_text(
        encoding="utf-8"
    )


def read_file(file_path: str) -> str:
    """
    Automatically detect the file type and
    extract its contents.
    """

    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(
            f"File not found: {file_path}"
        )

    extension = path.suffix.lower()

    if extension == ".pdf":
        return read_pdf(file_path)

    if extension == ".docx":
        return read_docx(file_path)

    if extension == ".pptx":
        return read_pptx(file_path)

    if extension == ".csv":
        return read_csv(file_path)

    if extension in [".xlsx", ".xls"]:
        return read_xlsx(file_path)

    if extension == ".txt":
        return read_txt(file_path)

    raise ValueError(
        f"Unsupported file format: {extension}"
    )